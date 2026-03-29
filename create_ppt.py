# -*- coding: utf-8 -*-
"""
Create PPT for Momentum Factor Reproduction
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define title layout
title_layout = prs.slide_layouts[6]  # Blank layout

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(title_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(title_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.space_after = Pt(12)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(title_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Table
    x = Inches(1)
    y = Inches(1.2)
    width = Inches(11.333)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(len(rows)+1, len(headers), x, y, width, height).table
    
    # Headers
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    
    # Data
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            table.cell(i+1, j).text = str(cell)
    
    return slide

# ===== SLIDES =====

# Slide 1: Title
add_title_slide(prs, 
    "成交量对动量因子的修正：日与夜的殊途同归",
    "东吴证券金融工程专题报告复现\n报告人：XXX\n日期：2025年XX月XX日")

# Slide 2: Contents
add_content_slide(prs, "目录", [
    "01 研究背景与核心发现",
    "02 数据获取与预处理",
    "03 因子构建方法",
    "04 实证结果分析",
    "05 稳健性检验",
    "06 结论与技术实现"
])

# Slide 3: Background
add_content_slide(prs, "01 研究背景与核心发现", [
    "传统动量因子存在收益不稳定、回撤大的问题",
    "核心发现：",
    "  • 日内：高换手率强化动量信号（锦上添花）",
    "  • 隔夜：昨日换手率影响收益方向（雪中送炭）",
    "新动量因子在稳定性、胜率上显著优于传统因子"
])

# Slide 4: Formulas
add_content_slide(prs, "02 核心公式", [
    "收益定义：",
    "  • 日内收益: r_t = close_t / open_t - 1",
    "  • 隔夜收益: g_t = open_t / close_{t-1} - 1",
    "",
    "因子合成：",
    "  • NEW_Intraday = -zscore(part1) + zscore(part5)",
    "  • NEW_Overnight = +zscore(part1) - zscore(part5)",
    "  • NEW_Momentum = zscore(NEW_Intraday) + zscore(NEW_Overnight)"
])

# Slide 5: Data
add_content_slide(prs, "03 数据说明", [
    "数据来源：A股日频行情数据",
    "时间范围：2010-2019年",
    "核心字段：",
    "  • OHLC价格",
    "  • 成交量 (volume)",
    "  • 换手率 (turnover_rate)",
    "  • 股票状态（ST、停牌）",
    "处理方式：剔除ST、停牌股票，处理一字涨跌停"
])

# Slide 6: Process
add_content_slide(prs, "04 因子构建流程", [
    "Step 1: 数据预处理",
    "  - 计算日内/隔夜收益、换手率",
    "",
    "Step 2: 传统因子",
    "  - OLD_Momentum / OLD_Intraday / OLD_Overnight",
    "",
    "Step 3: 局部因子",
    "  - 按换手率排序分组，计算组内均值",
    "",
    "Step 4: 新因子合成"
])

# Slide 7: IC Results
add_table_slide(prs, "05 IC分析结果",
    ["因子", "IC均值", "年化ICIR", "IC T统计量"],
    [
        ["OLD_Momentum", "-0.030", "1.04", "2.08"],
        ["OLD_Intraday", "-0.028", "0.98", "1.96"],
        ["OLD_Overnight", "-0.032", "1.12", "2.24"],
        ["NEW_Intraday", "-0.041", "1.83", "3.66"],
        ["NEW_Overnight", "-0.049", "2.31", "4.62"],
        ["NEW_Momentum", "-0.055", "3.04", "6.08"]
    ])

# Slide 8: Backtest Results
add_table_slide(prs, "06 5分组回测结果",
    ["指标", "传统动量", "新动量", "改善"],
    [
        ["年化收益率", "19.71%", "18.65%", "-1.06%"],
        ["信息比率", "1.04", "2.89", "+1.85"],
        ["月度胜率", "66.15%", "86.15%", "+20%"],
        ["最大回撤", "15.89%", "6.33%", "-9.56%"]
    ])

# Slide 9: Robustness
add_content_slide(prs, "07 稳健性检验", [
    "改变回看天数：N = 40、60日",
    "扩展样本起点：2010年",
    "不同样本空间：",
    "  • 沪深300成分股",
    "  • 中证500成分股",
    "Barra中性化：",
    "  • 10个Barra风格因子",
    "  • 28个申万一级行业"
])

# Slide 10: Conclusion
add_content_slide(prs, "08 结论", [
    "1. 成交量信息可有效增强动量因子",
    "2. 新动量因子稳定性显著提升",
    "3. 月度胜率从66%提高至86%",
    "4. 最大回撤从16%降低至6%",
    "5. ICIR从1.04提升至3.04"
])

# Slide 11: Tech Stack
add_content_slide(prs, "09 技术实现", [
    "编程语言：Python 3.10+",
    "核心库：",
    "  • pandas, numpy - 数据处理",
    "  • scipy - 统计分析",
    "  • numba - 计算加速",
    "  • matplotlib - 可视化",
    "数据来源：Tushare / AKShare",
    "代码管理：Git + GitHub"
])

# Slide 12: Thank you
add_title_slide(prs, "谢谢！", "欢迎提问与讨论")

# Save
output = "D:/3YearSpring/FBDQA/东吴/动量因子复现报告.pptx"
prs.save(output)
print(f"Saved to {output}")
